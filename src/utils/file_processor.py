# File Processing Utilities
import os
import asyncio
from typing import List, Dict, Any, Tuple
from fastapi import UploadFile
from PyPDF2 import PdfReader
from pptx import Presentation
import pandas as pd
import csv
from io import BytesIO, StringIO
import mimetypes

class FileProcessor:
    """Handles file processing and content extraction for multiple file types."""
    
    @staticmethod
    async def extract_text_from_upload(file: UploadFile) -> Tuple[str, Dict[str, Any]]:
        """Extract text from uploaded file and return content with metadata."""
        filename = (file.filename or "").lower()
        content = await file.read()
        
        metadata = {
            "filename": file.filename,
            "size_bytes": len(content),
            "content_type": file.content_type,
            "file_extension": os.path.splitext(filename)[1] if filename else ""
        }
        
        try:
            # Text files
            if filename.endswith((".txt", ".md")):
                text_content = FileProcessor._extract_text_content(content)
                metadata["processing_method"] = "text_decode"
                return text_content, metadata
            
            # CSV files
            elif filename.endswith(".csv"):
                text_content = FileProcessor._extract_csv_content(content)
                metadata["processing_method"] = "csv_parse"
                return text_content, metadata
            
            # Excel files
            elif filename.endswith((".xlsx", ".xls")):
                text_content = FileProcessor._extract_excel_content(content)
                metadata["processing_method"] = "excel_parse"
                return text_content, metadata
            
            # PDF files
            elif filename.endswith(".pdf"):
                text_content = FileProcessor._extract_pdf_content(content)
                metadata["processing_method"] = "pdf_extract"
                return text_content, metadata
            
            # PowerPoint files
            elif filename.endswith(".pptx"):
                text_content = FileProcessor._extract_pptx_content(content)
                metadata["processing_method"] = "pptx_extract"
                return text_content, metadata
            
            # JSON files (for web scraping results)
            elif filename.endswith(".json"):
                text_content = FileProcessor._extract_json_content(content)
                metadata["processing_method"] = "json_parse"
                return text_content, metadata
            
            else:
                # Try to detect content type
                mime_type, _ = mimetypes.guess_type(filename)
                if mime_type and mime_type.startswith("text/"):
                    text_content = FileProcessor._extract_text_content(content)
                    metadata["processing_method"] = "text_fallback"
                    return text_content, metadata
                else:
                    raise ValueError(f"Unsupported file type: {filename}")
                    
        except Exception as e:
            metadata["error"] = str(e)
            return f"Error processing file {filename}: {str(e)}", metadata
    
    @staticmethod
    def _extract_text_content(content: bytes) -> str:
        """Extract text from text files."""
        try:
            return content.decode("utf-8", errors="ignore")
        except Exception:
            return content.decode("latin-1", errors="ignore")
    
    @staticmethod
    def _extract_csv_content(content: bytes) -> str:
        """Extract and format CSV content."""
        try:
            text_content = content.decode("utf-8", errors="ignore")
        except UnicodeDecodeError:
            text_content = content.decode("latin-1", errors="ignore")
        
        csv_reader = csv.reader(StringIO(text_content))
        rows = list(csv_reader)
        
        if not rows:
            return "Empty CSV file"
        
        headers = rows[0] if rows else []
        formatted_text = [
            f"CSV Data Analysis:",
            f"Rows: {len(rows)}, Columns: {len(headers)}",
            f"Headers: {', '.join(headers)}",
            "\nData Sample:"
        ]
        
        # Add sample rows
        for i, row in enumerate(rows[1:21], 1):  # First 20 rows
            formatted_text.append(f"Row {i}: {', '.join(str(cell) for cell in row)}")
        
        if len(rows) > 21:
            formatted_text.append(f"... and {len(rows) - 21} more rows")
        
        return "\n".join(formatted_text)
    
    @staticmethod
    def _extract_excel_content(content: bytes) -> str:
        """Extract and format Excel content."""
        df_dict = pd.read_excel(BytesIO(content), sheet_name=None)
        
        formatted_text = [f"Excel Analysis - {len(df_dict)} sheet(s):"]
        
        for sheet_name, df in df_dict.items():
            formatted_text.append(f"\n--- Sheet: {sheet_name} ---")
            formatted_text.append(f"Dimensions: {df.shape[0]} rows × {df.shape[1]} columns")
            
            if not df.empty:
                formatted_text.append(f"Columns: {', '.join(df.columns.astype(str))}")
                formatted_text.append("\nSample Data:")
                
                for idx, (_, row) in enumerate(df.head(10).iterrows()):
                    row_text = ", ".join(f"{col}: {val}" for col, val in row.items() if pd.notna(val))
                    formatted_text.append(f"  {idx + 1}. {row_text}")
                
                if len(df) > 10:
                    formatted_text.append(f"  ... and {len(df) - 10} more rows")
        
        return "\n".join(formatted_text)
    
    @staticmethod
    def _extract_pdf_content(content: bytes) -> str:
        """Extract text from PDF files."""
        reader = PdfReader(BytesIO(content))
        text_parts = []
        
        for page_num, page in enumerate(reader.pages, 1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(f"--- Page {page_num} ---\n{page_text}")
        
        return "\n\n".join(text_parts) if text_parts else "No extractable text found in PDF"
    
    @staticmethod
    def _extract_pptx_content(content: bytes) -> str:
        """Extract text from PowerPoint files."""
        prs = Presentation(BytesIO(content))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        
        return "\n\n".join(text_parts) if text_parts else "No extractable text found in presentation"
    
    @staticmethod
    def _extract_json_content(content: bytes) -> str:
        """Extract and format JSON content."""
        import json
        try:
            data = json.loads(content.decode("utf-8"))
            return f"JSON Data:\n{json.dumps(data, indent=2, ensure_ascii=False)}"
        except Exception as e:
            return f"Error parsing JSON: {str(e)}"
    
    @staticmethod
    def categorize_file(filename: str, content_type: str = None) -> str:
        """Categorize file into data source types."""
        filename_lower = filename.lower()
        
        # Pitch deck materials
        if any(keyword in filename_lower for keyword in ["pitch", "deck", "presentation", "slides"]):
            return "pitch_deck"
        
        # Data room documents
        elif any(keyword in filename_lower for keyword in ["financial", "traction", "metrics", "kpi", "revenue", "onepager", "one-pager"]):
            return "data_room"
        
        # Web content
        elif any(keyword in filename_lower for keyword in ["website", "landing", "blog", "faq", "web", "content"]):
            return "web_content"
        
        # Interaction signals
        elif any(keyword in filename_lower for keyword in ["call", "recording", "transcript", "interview", "questionnaire", "feedback"]):
            return "interaction"
        
        # Default categorization by file type
        elif filename_lower.endswith((".pdf", ".pptx", ".key")):
            return "pitch_deck"
        elif filename_lower.endswith((".xlsx", ".xls", ".csv")):
            return "data_room"
        elif filename_lower.endswith((".txt", ".md", ".json")):
            return "web_content"
        else:
            return "general"
